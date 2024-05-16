import os

import json
# importing required modules
from datetime import datetime
from math import floor

from ShiftArchi import build_archi_xml
from ShiftJson import dump_csv
from pptx import Presentation
from MsoShapeList import MSO_AUTO_SHAPE_TYPE_DICT, MSO_SHAPE_TYPE_DICT
# from pptx.enum.shapes import MSO_SHAPE
import numpy as np
import cv2
# from pptx.enum.shapes import MSO_SHAPE_TYPE
from spire.presentation.common import *
from spire.presentation import Presentation as ShapePresentation, ISmartArt
from PIL import Image

from pathlib import Path
import logging

DOWNLOADS_PATH = str(Path.home()/"Downloads")




# def find_connected_shapes(slide_number, slide, text_runs):
# # Iterate through all shapes on the slide
#     for shape in slide.shapes:
#         if shape.shape_type == MSO_SHAPE_TYPE.LINE:
#             # This is a connector
#             connector = shape
#             # Assuming the connector has a start and end, check positions (this is hypothetical)
#             # This part is not supported by python-pptx and is here for illustrative purposes
#             start_pos = [connector.begin_x, connector.begin_y]
#             start_pos_inch = [connector.begin_x.inches, connector.begin_y.inches]
#             end_pos = [connector.end_x, connector.end_y]
#             end_pos_inch = [connector.end_x.inches, connector.end_y.inches]
#             start_shape = None
#             end_shape = None
#             # Find shapes that match the start and end positions
#             for other_shape in slide.shapes:
#                 if other_shape.shape_type != MSO_SHAPE_TYPE.LINE:
#                     if is_near(other_shape, start_pos, start_pos_inch):
#                         start_shape = other_shape
#                     if is_near(other_shape, end_pos, end_pos_inch):
#                         end_shape = other_shape
#             return start_shape, end_shape
#
# def is_near(shape, position, position_inch):
#     # Dummy function to check if a position is within a shape's bounds
#     # This would need to calculate based on the actual position data
#     margin = 100
#     check = False
#     top = shape.top
#     left = shape.left
#     width = shape.width
#     height = shape.height
#     x1, y1, w, h = (left, top, width, height)
#     x2, y2 = x1 + w, y1 + h
#     x, y = position
#     test_x1 = x1 - margin
#     test_x2 = x2 + margin
#     test_y1 = y1 - margin
#     test_y2 = y2 + margin
#     y1_inch = shape.top.inches
#     x1_inch = shape.left.inches
#     width_inch = shape.width.inches
#     height_inch = shape.height.inches
#     x2_inch = x1_inch + width_inch
#     y2_inch = y1_inch + height_inch
#     horizontal = position_inch[0]
#     vertical = position_inch[1]
#     print (f" {horizontal:.3}, {vertical:.3}")
#     if test_x1 < x < test_x2:
#         if test_y1 < y < test_y2:
#             check = True
#     return check


def is_near_dict(other, position):
    margin = 10000
    check = False

    top = other['top']
    left = other['left']
    width = other['width']
    height = other['height']

    x1, y1, w, h = (left, top, width, height)
    x2, y2 = x1 + w, y1 + h
    x, y = position
    test_x1 = x1 - margin
    test_x2 = x2 + margin
    test_y1 = y1 - margin
    test_y2 = y2 + margin
    # em per inch is 914400
    print(f" LINE: position: {x / 914400:.3}, {y / 914400:.3}")
    print(f" {other ['text']}: x check: {test_x1/914400:.3} < {x/914400:.3} < {test_x2/914400:.3}")
    print(f"                   y check: {test_y1/914400:.3} < {y/914400:.3} <{test_y2/914400:.3}")
    # print(f" (horizontal:.3}, {vertical:.3}")
    if test_x1 < x < test_x2:
        if test_y1 < y < test_y2:
            check = True
    return check


def find_connected_shapes_json(text_runs):
    for run in text_runs:
        if run['type'] == "LINE":
            start_shape_id = None
            end_shape_id = None
            print(f" {run['begin']}, {run['end']}")
            for other_run in text_runs:
                if other_run['type'] != "LINE":
                    if is_near_dict(other_run, run['begin']):
                        start_shape_id = other_run['id']
                    if is_near_dict(other_run, run['end']):
                        end_shape_id = other_run['id']

            run['start_id'] = start_shape_id
            run['end_id'] = end_shape_id

    return text_runs


def convert_slide_json_to_connect_struct(model_name, slide_list_dict):
    elements = []
    scaling_factor = 80 / 914400 # creates an 80 pixel per inch scaling.  (914400 em is one inch)
    relationships = []
    organizations = [{'label': 'Business', 'identifier': []},
                     {'label': 'Application', 'identifier': []},
                     {'label': 'Relations', 'identifier': []},
                     {'label': 'Views', 'identifier': []}]
    diagram_view_components = []
    for item in slide_list_dict:
        if item['type'] == "AUTOSHAPE: RECTANGLE":
            # 1) add an element type to elements
            # elements.append({'type': "Application Component", 'identifier': "app-1", 'name': 'Online Bookstore',
            #            'documentation': 'A web application allowing users to buy books and manage their purchases. '})
            # for each of the elements of the same type (the text in the block identifies the type), the element is
            # for example:  item_id = f"id-{item['id']}"
            item_id = f"id-instance{item['id']}"

            identifier_id =  f"app-{item['id']}"
            # build the element component
            # {'type': "BusinessActor",  << This will depend on the item in brackets
            #                          'identifier': "actor-1",
            #                          'name': "Customer",
            #                          'documentation': "Any person who purchases books from the Online Bookstore."}
            elements.append({'type': "ApplicationComponent",
                             'identifier': identifier_id,
                             'name': item['name'],
                             'documentation': item.get('description', 'none')})
            # 2) and add a diagram entry. This is the actual placement of the component on the page.
            diagram_view_components.append({'type': 'element',
                                            'id': item_id,
                                            'ref': identifier_id,
                                            'box': {'x': round(item['left'] * scaling_factor),
                                                    'y': round(item['top'] * scaling_factor),
                                                    'w': round(item['width'] * scaling_factor),
                                                    'h': round(item['height'] * scaling_factor)}
                                            })
            # 3) collect an organizations entry (This is like the catalog)
            # {'label': 'Application', identifier': ["app-1", "ext-1"]},
            for org_item in organizations:
                if org_item['label'] == 'Application':
                    org_item['identifier'].append(identifier_id)

        elif item["type"] == "LINE":
            # 1) add a relationship
            # {'type': "Association",
            #  'identifier': "rel-1",
            #  'source': "app-1",
            #  'target': "actor-1",
            #  'documentation': "Allows customers to browse, search, and buy books."}


            ref_id = f"rel-{item['id']}"
            identifier_id = f"id-connector{item['id']}"
            source_rel = f"app-{item['start_shape_id']}"
            target_rel = f"app-{item['end_shape_id']}"
            source_view_id = f"id-instance{item['start_shape_id']}"
            target_view_id = f"id-instance{item['end_shape_id']}"
            relationships.append({'type': "Association",
                                  'identifier': ref_id,
                                  'source': source_rel,
                                  'target': target_rel,
                                  'documentation': item.get('documentation', 'none')})
            # 2) and add a diagram entry. This is the actual placement of the component on the page.
            # {'type': 'relationship',
            #  'id': "id-connector1"
            #  'ref': "rel-2"
            #  'source': "id-instance1"
            #  'target': id-instance3"} ]

            diagram_view_components.append({'type': 'relationship',
                                            'id': identifier_id,
                                            'ref': ref_id,
                                            'source': source_view_id,
                                            'target': target_view_id})

            # 3) collect an organizations entry
            # {'label': 'Relations', 'identifier': ["rel-1", "rel-2", "rel-3"]}
            for org_item in organizations:
                if org_item['label'] == 'Relations':
                    org_item['identifier'].append(ref_id)

    # build the diagram view
    file_id = 'id-1234'
    view_name = "Context View"
    view_instance_name = "view-1"
    for org_item in organizations:
        if org_item['label'] == 'Views':
            org_item['identifier'].append(view_instance_name)
    # sort the diagram view components to put elements and relationships together
    dvc = sorted(diagram_view_components, key=lambda x: x['type'])
    s = build_archi_xml(file_id, model_name, elements, relationships, organizations,
                        view_instance_name, dvc, view_name)
    return s


def image_text_extract(png_file):
    global reader
    # text = pytesseract.image_to_string (png_file)
    results = reader.readtext(png_file)
    all_text = []
    for (bbox, text, prob) in results:
        bbox = [list(map(round, sublist)) for sublist in bbox]
        top = bbox[0][1]
        left = bbox[0][0]
        width = bbox[1][0]
        height = bbox[3][1] - bbox[0][1]
        t = {'image_bbox': {'left': left, 'top': top, 'width': width, 'height': height},
             'text': text,
             'prob': round(prob, 3)}
        all_text.append(t)
    return all_text


def on_line_trackbar(detection_threshold):
    global tracker_image, gray_img
    img = tracker_image.copy()
    # Edge detection
    hist1 = 50
    app_size = 7
    edges = cv2.Canny(gray_img, hist1, 150, apertureSize=app_size, L2gradient=True)
    # detection_threshold = 50
    # Detect lines
    linesP = cv2.HoughLinesP(edges, 1, np.pi / 180, detection_threshold, minLineLength=100, maxLineGap=10)

    # Draw lines
    for line in linesP:
        x1, y1, x2, y2 = line[0]
        cv2.line(img, (x1, y1), (x2, y2), (0, 255, 0), 2)
    cv2.imshow('lines', img)


def image_line_extract_tracker(png_file):
    global tracker_image, gray_img # , rotated_image
    # Read the image
    pre_image = cv2.imread(png_file)
    # Get image dimensions
    (h, w) = pre_image.shape[:2]
    center = (w // 2, h // 2)

    # Rotate 10 degrees clockwise
    angle = -10
    scale = 1.0

    # Get the rotation matrix
    M = cv2.getRotationMatrix2D(center, angle, scale)

    # Perform the rotation
    tracker_image = cv2.warpAffine(pre_image, M, (w, h))
    # Optionally, you can use cv2.imshow to see the rotated image
    # cv2.imshow('Rotated Image', rotated_image)

    gray_img = cv2.cvtColor(tracker_image, cv2.COLOR_BGR2GRAY)
    gray_img = cv2.equalizeHist(gray_img)
    # gray_img = cv2.GaussianBlur(gray_img, (3, 3), 0)

    # Create a window and a trackbar
    cv2.namedWindow('lines')
    cv2.createTrackbar('Threshold', 'lines', 50, 255, on_line_trackbar)

    # Initialize with a default value
    on_line_trackbar(50)

    # Display the window until a key is pressed
    cv2.waitKey(0)
    cv2.destroyAllWindows()
    return 1


def process_shape(output_path, prs_shaper, slide_number, shape):

    # No text frame
    t_runs = []
    a_t_runs = []
    if shape.has_text_frame:
        if MSO_SHAPE_TYPE_DICT[shape.shape_type]['name'] == "AUTO_SHAPE":
            shape_type = f"AUTOSHAPE: {MSO_AUTO_SHAPE_TYPE_DICT[shape.shape_type]['name']}"
        else:
            shape_type = MSO_SHAPE_TYPE_DICT[shape.shape_type]['name']
        logger.info(f"Processing {shape_type} on slide {slide_number + 1} at "
                    f"left{round(shape.left.inches, 3)}, top: {round(shape.top.inches, 3)}")
        left = shape.left
        top = shape.top
        # width = shape.width
        height = shape.height
        rows = len(shape.text_frame.paragraphs)
        if rows > 0:
            y_incr = height / rows
            for i, paragraph in enumerate(shape.text_frame.paragraphs):
                for run in paragraph.runs:
                    t = {'slide': slide_number + 1,
                         'id': shape.shape_id,
                         'type': shape_type,
                         'name': shape.name,
                         'left': left,
                         'top': floor(top + y_incr * i),
                         'width': shape.width,
                         'height': shape.height,
                         'text': run.text}
                    t_runs.append(t)
                    a_t_runs.append(t)
        return t_runs, a_t_runs
    else:
        if shape.shape_type is None:
            shape_slide = prs_shaper.Slides[slide_number]
            for smart_art_shape in shape_slide.Shapes:
                if isinstance(smart_art_shape, ISmartArt):
                    # Extract text from the SmartArt shapes and append the text to the list
                    for n, node in enumerate(smart_art_shape.Nodes):
                        t = {'slide': slide_number + 1,
                             'id': shape.shape_id,
                             "type": "SmartArt",
                             'name': shape.name,
                             'left': shape.left,
                             'top': shape.top,
                             # left,
                             # floor (top + y_incr * 1),
                             'width': shape.width,
                             'height': shape.height,
                             'node': n,
                             'text': node.TextFrame.Text}
                        t_runs.append(t)
                        a_t_runs.append(t)
            return t_runs, a_t_runs
        # Some other shape without a text frame
        shape_name = MSO_SHAPE_TYPE_DICT[shape.shape_type].get('name', "NONE")
        if shape_name == "TABLE":
            # process the table here
            logger.info(f"Processing {MSO_SHAPE_TYPE_DICT[shape.shape_type]['name']} on slide {slide_number + 1} "
                        f"at left: {round(shape.left.inches, 3)}, top: {round(shape.top.inches, 3)}")
            num_rows = len(shape.table.rows)
            num_columns = len(shape.table.columns)
            for r in range(num_rows):
                for c in range(num_columns):
                    t = {'slide': slide_number + 1,
                         'id': shape.shape_id,
                         "type": MSO_SHAPE_TYPE_DICT[shape.shape_type]['name'],
                         'left': shape.left, # left,
                         'top': shape.top,  # floor (top + y_incr + i),
                         'width': shape.width,
                         'height': shape.height,
                         'row': r,
                         'col': c,
                         'text': shape.table.cell(r, c).text}
                    t_runs.append(t)
                    a_t_runs.append(t)
            return t_runs, a_t_runs
        elif shape_name == "PICTURE":
            logger.info(f"Processing {shape_name} on slide {slide_number + 1} at left: {round(shape.left.inches, 3)}, top: {round(shape.top.inches, 3)}")
            image_bytes = shape.image.blob
            if image_bytes[0:10] == b'\xff\x68\xff\xe0\x00\x14JFIF':
                image_extension = 'jpg'
            else:
                image_extension = shape.image.ext
            if image_extension.lower() not in ["bmp", "gif", "jpg", "png", "tiff", "wmf"]:
                logger.info(f"Need special processing of the PICTURE of type {image_extension}, skipping")
                return t_runs, a_t_runs
            image_name = f"{output_path}/image_{shape.shape_id}.{image_extension}"
            with open(image_name, 'wb') as f:
                f.write(image_bytes)
            image_text_array = []
            img_height = 0
            img_width = 0
            try:
                # Image text extraction is tricky and may fail.
                if image_extension.lower() == "gif":
                    new_image_name = "{output_path}/image_{shape.shape_id}.png"
                    logger.info(f"Converting {image_name} into {new_image_name}")
                    with Image.open(image_name) as img:
                        img = img.convert("RGBA")
                        img.save(new_image_name, "PNG")
                elif image_extension.lower() == "wmf":
                    new_image_name = f"{output_path} / image_{shape.shape_id}.png"
                    logger.info(f"Converting {image_name} into {new_image_name}")
                    with Image.open(image_name) as img:
                        img = img.convent("RGBA")
                        img.save(new_image_name, "PNG")
                else:
                    new_image_name = image_name
                # Extract the text in this image
                image_text_array = image_text_extract(new_image_name)
                # Extract the shapes in this image
                # image_shape_array = image_shape_extract(new_image_name)
                # image_shape_array = image_shape_extract_tracker(new_image_name)
                image_line_array = image_line_extract_tracker(new_image_name)
                with Image.open(new_image_name) as img:
                    img_height = img.height
                    img_width = img.width
            except Exception as e:
                logger.error(f" Conversion failed for {image_name}, {e}")
            t = {'slide': slide_number + 1,
                 'id': shape.shape_id,
                 "type": "PICTURE",
                 'left': shape.left,
                 'top': shape.top,
                 'width': shape.width,
                 'height': shape.height,
                 'image_width': img_width,
                 'image_height': img_height,
                 'image_name': image_name,
                 'text': image_text_array}
            t_runs.append(t)
            a_t_runs.append(t)
            return t_runs, a_t_runs
        elif shape_name == "GROUP":
            logger.info(f"Processing a group of {len(shape.shapes)} items on slide {slide_number + 1}")
            for grouped_shape in shape.shapes:
                t_r, a_t_r = process_shape(output_path, prs_shaper, slide_number, grouped_shape)
                t_runs.extend(t_r)
                a_t_runs.extend(a_t_r)
            return t_runs, a_t_runs
        elif shape_name == "LINE":
            start_connector = shape.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}stCxn')
            if start_connector is not None:
                start_shape_id = int(start_connector.get('id'))
                start_shape_index = int(start_connector.get('idx'))
            else:
                start_shape_id = None
                start_shape_index = None
            end_connector = shape.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}endCxn')
            if end_connector is not None:
                end_shape_id = int(end_connector.get('id'))
                end_shape_index = int(end_connector.get('idx'))
            else:
                end_shape_id = None
                end_shape_index = None
            t = {'slide': slide_number + 1,
                 'id': shape.shape_id,
                 'name': shape.name,
                 "type": "LINE",
                 'left': shape.left,  # left,
                 'top': shape.top,  # floor (top + y_incr * i),
                 'begin': [shape.begin_x, shape.begin_y],
                 'end': [shape.end_x, shape.end_y],
                 'start_shape_id': start_shape_id,
                 'start_shape_index': start_shape_index,
                 'end_shape_id': end_shape_id,
                 'end_shape_index': end_shape_index,
                 }
            logger.debug(f"{shape.name} begin: {shape.begin_x / 914400 - 13.333 / 2:.3}, "
                         f"{shape.begin_y / 914400 - 7.5 / 2:.3}, "
                         f"end: {shape.end_x / 914400 - 13.333 / 2:.3}, {shape.end_y / 914400 - 7.5 / 2:.3}")
            t_runs.append(t)
            a_t_runs.append(t)
            return t_runs, a_t_runs
        else:
            # LINE, FREEFORM, EMBEDDED_OLE_OBJECT, AUTO_SHAPE, TEXT_BOX, MEDIA
            logger.warning(f"Skipping {shape_name} on slide {slide_number + 1} at left: {round(shape.left.inches, 3)},"
                           f" top: {round(shape.top.inches, 3)}")
            return t_runs, a_t_runs


def ppt_extraction(path_to_presentation):
    prs = Presentation(path_to_presentation)
    prs_shaper = ShapePresentation()
    prs_shaper.LoadFromFile(path_to_presentation)
    base_name = os.path.splitext(os.path.basename(path_to_presentation))[0]
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    prefix = datetime.now().strftime("Y%Ym%md%d-H%HM%MS%S")
    extract_time = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
    output_path = f"./pptx/scan/{prefix}/"
    path_exists = os.path.exists(output_path)
    if not path_exists:
        # Create a new directory because it does not exist
        os.makedirs(output_path)
    logging.basicConfig(filename=f'{output_path}conversion.log', encoding='utf-8', level=logging.INFO)
    logging.info(f"Processing the power point file {path_to_presentation}")

    meta_data = {"file": path_to_presentation,
                 "slide count": len(prs.slides),
                 "author": prs.core_properties.author,
                 "created on": f" {prs.core_properties.created} UTC",
                 "modified on": f" {prs.core_properties.modified} UTC",
                 "extracted on": f" {extract_time} UTC"}
    with open(f"{output_path}/meta-data.json", 'w') as f:
        f.write(json.dumps(meta_data, indent=4))
    all_text_list = []
    for slide_number, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_number + 1} of {len (prs.slides)}")
        # getting a specific page from the per file
        output_path = f"./pptx/scan/{prefix}/slide-{slide_number + 1}"
        path_exists = os.path.exists(output_path)
        if not path_exists:
            # Create a new directory because it does not exist
            os.makedirs(output_path)
        text_runs = []
        for shape in slide.shapes:
            t_runs, a_t_runs = process_shape(output_path, prs_shaper, slide_number, shape)
            text_runs.extend(t_runs)
            all_text_list.extend(a_t_runs)
        # new text runs find connected shapes.json(text_runs)
        # if new_text_runs:
        #     print("Start Shape:", new text_runs name if new text runs [0] else "None")
        #     print("End Shape:", new text runs [1] name if new text runs [1] else "None")
        with open(f"{output_path}/text_runs.json", 'a', encoding="utf-8") as f:
            f.write(json.dumps({'page text': text_runs}, indent=4))
            # for t in text_runs:
            #     f.write(f"{t['slide']}|{t['x']}|{t['y']){t['text']}\n")
            #     all_text_list.append(t)
            #     f.write(f*(bytes (t, 'utf-8').decode('unicode_escape')}\n*)
        xml_string = convert_slide_json_to_connect_struct(f"{base_name} slide-{slide_number + 1}", text_runs)
        with open(f"{output_path}/slide_arch.xml", 'a', encoding="utf-8") as f:
            f.write(xml_string)

    output_path = f"./pptx/scan/{prefix}"
    path_exists = os.path.exists(output_path)
    if not path_exists:
        # Create a new directory because it does not exist
        os.makedirs(output_path)
    # Zero Width String
    # joined = '\n'.join(json.dumps (all_text_list))
    slide_text = json.dumps({'deck_text': all_text_list}, indent=4)
    with open(f"{output_path}/joined.json", 'a', encoding="utf-8") as f:
        f.write(slide_text)

    dump_csv(all_text_list, f"{output_path}/joined.csv")

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)


if __name__ == "__main__":
    # convert_pdf_to_png('./incoming/ems-color-system-diagram.pdf')
    run_directory = os.getcwd()
    # slide_path = "US Electric GIS_Strategy_v01_debrief_11_21_22_summary.pptx
    # slide_path = "PI Facilitation Pack March 2024 Plan. Work Execution v2.pptx"
    slide_path = "./incoming/connection.pptx"
    print(f"Processing '{slide_path}'")
    ppt_extraction(slide_path)
